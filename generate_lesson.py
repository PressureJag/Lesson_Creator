#!/usr/bin/env python3
"""
Lesson slide-deck generator for Outwood Grange Academies Trust.

Usage — interactive (recommended):
  python3 generate_lesson.py

Usage — direct PDF paths:
  python3 generate_lesson.py \\
    --summary  "Examples/SOW/Algebra 1 - Core Plus Summary of Intent.pdf" \\
    --methods  "Examples/SOW/Algebra 1 - Consistent Methodology.pdf" \\
    --topic    "Algebra 1 Core Plus" \\
    --year     "Year 9" \\
    --pitch    "Core Plus" \\
    --output   "Output/Algebra_1.pptx"
"""

import argparse
import os
import sys

from pptx import Presentation

from generator import sow_parser, content_gen, slide_builder, diagram_gen, planner
from generator import theme as T
from generator.notification import header, notify, prompt_choice
from generator import gdrive
from generator import module_selector, class_profile as cp


# ── Diagram builder ──────────────────────────────────────────────────────────

def build_diagram(obj_text: str, diagram_override=None):
    """Generate an appropriate diagram image, or return None."""
    from generator import content_gen as cg
    dtype = diagram_override or cg.select_diagram_type(obj_text)

    if dtype == "percentage_bar":
        return diagram_gen.percentage_bar(
            whole=240, part_pct=30,
            whole_label="100%", part_label="30%",
            title="Percentage model"
        )
    if dtype == "number_line":
        return diagram_gen.number_line(
            lo=-6, hi=8,
            intervals=[(-2, 5, True, False)],
            title="Number line"
        )
    if dtype == "area_grid":
        return diagram_gen.area_grid(
            rows=["x", "+3"],
            cols=["x", "+2"],
            cells=[["x²", "+2x"], ["+3x", "+6"]],
            title="Expansion grid: (x+3)(x+2)"
        )
    if dtype == "angle":
        return diagram_gen.angle_diagram("alternate")
    if dtype == "bar_model":
        return diagram_gen.bar_model(
            parts=[(1, "#D4C5E2", "a"), (2, "#5B2C8D", "2a")],
            total_label="3a"
        )
    return None


# ── Slide deck assembly ──────────────────────────────────────────────────────

def build_deck(sow: dict, methods: dict, topic_name: str, output_path: str) -> None:
    import json as _json

    prs = Presentation()
    prs.slide_width  = T.SLIDE_WIDTH
    prs.slide_height = T.SLIDE_HEIGHT

    objectives = sow["objectives"]
    prior      = sow["prior_knowledge"]
    misconcs   = sow["misconceptions"]
    vocab      = sow["vocabulary"]
    extend     = sow["extend"]
    pd_        = sow["personal_development"]
    hours      = sow["time_hours"]
    big_q      = sow["big_question"]

    total = len(objectives)
    methods_pages = list(methods.values())

    # Content log — saved as JSON alongside the PPTX for review/editing
    content_log = {"topic": topic_name, "objectives": []}

    header(
        f"BUILDING DECK  —  {topic_name}",
        f"{total} objectives · {total * 15} slides estimated",
    )
    notify(f"Objectives parsed: {total}", "info")
    notify(f"Vocabulary terms:  {len(vocab)}", "info")
    notify(f"Misconceptions:    {len(misconcs)}", "info")

    # ── Intro slides ─────────────────────────────────────────
    notify("Building title slide …", "progress")
    slide_builder.make_title_slide(prs, topic_name, big_q, hours)

    notify("Building overview (Big Picture) slide …", "progress")
    slide_builder.make_overview_slide(
        prs, topic_name,
        sow.get("prior_knowledge", []),
        objectives,
        sow.get("future_knowledge", []),
    )

    notify("Building prior knowledge slide …", "progress")
    slide_builder.make_prior_knowledge(prs, topic_name, prior)

    if vocab:
        notify("Building vocabulary slide …", "progress")
        slide_builder.make_vocabulary(prs, topic_name, vocab)

    # ── Objective blocks ─────────────────────────────────────
    skip_next = False
    abort     = False
    idx       = 0

    for idx, obj in enumerate(objectives, start=1):

        if skip_next:
            notify(f"Skipping objective {idx}/{total} (user request)", "warning")
            skip_next = False
            continue

        # ── Planning choices ──────────────────────────────────
        obj_short = obj[:70] + ("…" if len(obj) > 70 else "")

        if idx - 1 < len(methods_pages):
            methods_text = methods_pages[idx - 1]
        else:
            methods_text = ""
            for key, val in methods.items():
                if any(w in key.lower() for w in obj.lower().split()[:4]):
                    methods_text = val
                    break

        plan_opts = planner.get_options(obj, topic_name, methods_text)
        plan_display = planner.to_prompt_list(plan_opts)

        choice_idx = prompt_choice(
            title=f"OBJECTIVE {idx}/{total}  —  {obj_short}",
            options=plan_display,
            default=0,
            timeout=120,
        )

        chosen = plan_opts[choice_idx]
        if chosen["skip"]:
            notify(f"Skipping objective {idx}/{total}", "warning")
            continue

        # ── Build slides for this objective ───────────────────
        notify(f"Objective {idx}/{total}: {obj[:55]} …", "section")

        obj_misconc = [misconcs[idx - 1]] if idx <= len(misconcs) else []

        # 1 & 2 — Starter + Answers (plus layout)
        notify("Generating retrieval starter …", "progress")
        retrieval = content_gen.generate_retrieval_questions(prior, obj, topic_name, vocab)
        slide_builder.make_starter_plus(prs, topic_name, retrieval["questions"])
        slide_builder.make_starter_plus(
            prs, topic_name,
            retrieval["questions"], retrieval["answers"]
        )

        # 3 — Learning Objective / Intro
        slide_builder.make_learning_intro(prs, topic_name, obj, idx)

        # 4, 5, 6 — I Do / We Do / You Do
        notify("Generating I Do / We Do / You Do sequence …", "progress")
        teaching = content_gen.generate_teaching_sequence(
            obj, topic_name, methods_text, vocab, obj_misconc
        )
        slide_builder.make_ido_slide(
            prs, topic_name,
            teaching["i_do"]["heading"],
            teaching["i_do"]["worked_example"],
            method_text=methods_text,
            notes=teaching["i_do"].get("notes", ""),
        )
        slide_builder.make_wedo_slide(
            prs, topic_name,
            teaching["we_do"]["heading"],
            teaching["we_do"]["question"],
            teaching["we_do"].get("steps"),
        )
        slide_builder.make_youdo_slide(
            prs, topic_name,
            teaching["you_do"]["heading"],
            teaching["you_do"]["question"],
            teaching["you_do"].get("answer", ""),
        )

        # 7–16 — Mini Whiteboard × 10
        notify("Generating mini whiteboard questions …", "progress")
        wb = content_gen.generate_mini_whiteboard_questions(obj, topic_name, methods_text)
        for wb_num, wb_q in enumerate(wb["questions"][:10], start=1):
            slide_builder.make_mini_whiteboard(prs, topic_name, wb_q, wb_num, 10)

        # 17 & 18 — Independent Practice + Answers
        notify("Generating independent practice …", "progress")
        indep = content_gen.generate_independent_practice(
            obj, topic_name, vocab, methods_text
        )
        slide_builder.make_independent_practice(prs, topic_name, indep["questions"])
        slide_builder.make_independent_answers(
            prs, topic_name, indep["questions"], indep["answers"]
        )

        # 19 — Plenary
        notify("Generating plenary …", "progress")
        plenary = content_gen.generate_plenary(obj, topic_name)
        slide_builder.make_plenary(
            prs, topic_name, obj,
            plenary["summary"], plenary["question"], plenary.get("answer", ""),
        )

        notify(f"Objective {idx}/{total} complete  ({len(prs.slides)} slides so far)", "success")

        # ── Log content for JSON export ───────────────────────
        content_log["objectives"].append({
            "index":             idx,
            "objective":         obj,
            "retrieval":         retrieval,
            "teaching_sequence": teaching,
            "mini_whiteboard":   wb,
            "independent":       indep,
            "plenary":           plenary,
            "misconception":     obj_misconc[0] if obj_misconc else None,
        })

        # ── Between-objective checkpoint ──────────────────────
        if idx < total:
            next_short = objectives[idx][:55] + ("…" if len(objectives[idx]) > 55 else "")
            is_last_remaining = (idx + 1 == total)
            skip_desc = (
                f"Omit the final objective and finish the deck"
                if is_last_remaining
                else f"Jump past objective {idx + 1} and continue from {idx + 2}"
            )
            cont_choice = prompt_choice(
                title=f"CHECKPOINT — after objective {idx}/{total}",
                options=[
                    ("Continue to next objective",
                     f"Next: {next_short}"),
                    ("Skip next objective", skip_desc),
                    ("Abort — save deck as-is",
                     "Stop here and save the partial deck to output + Google Drive"),
                ],
                default=0,
                timeout=120,
            )
            if cont_choice == 1:
                skip_next = True
            elif cont_choice == 2:
                notify("Aborting — saving partial deck …", "warning")
                abort = True
                break

    if abort:
        notify(f"Generation aborted after {idx} objective(s) — saving partial deck", "warning")

    # ── Extension ─────────────────────────────────────────────
    if not abort and extend:
        notify("Building extension slide …", "progress")
        slide_builder.make_extension(prs, topic_name, "Going Further", extend)

    # ── Save PPTX ─────────────────────────────────────────────
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)

    # ── Save content JSON (for review / refinement) ───────────
    content_path = output_path.replace(".pptx", "_content.json")
    with open(content_path, "w", encoding="utf-8") as f:
        _json.dump(content_log, f, indent=2, ensure_ascii=False)

    slide_count = len(prs.slides)
    header("GENERATION COMPLETE", f"{slide_count} slides  →  {output_path}", color="\033[32m")
    notify(f"Saved: {output_path}", "success")
    notify(f"Content JSON: {content_path}", "info")

    # ── Google Drive upload ────────────────────────────────────
    notify("Uploading to Google Drive …", "progress")
    result = gdrive.upload(output_path)
    if result["success"]:
        notify(f"Google Drive: {result['uploaded'][0]}", "success")
    else:
        notify(f"Google Drive upload failed: {result['error']}", "warning")

    # ── Final options ─────────────────────────────────────────
    final_choice = prompt_choice(
        title="WHAT WOULD YOU LIKE TO DO NEXT?",
        options=[
            ("Done",
             "Exit the generator"),
            ("Open output folder in Finder",
             f"Open {os.path.dirname(os.path.abspath(output_path))}"),
            ("Open Google Drive folder in Finder",
             result["folder"] if result["success"] else "Google Drive not available"),
        ],
        default=0,
        timeout=120,
    )

    if final_choice == 1:
        folder = os.path.dirname(os.path.abspath(output_path))
        os.system(f'open "{folder}"')
    elif final_choice == 2 and result["success"]:
        os.system(f'open "{result["folder"]}"')


# ── Startup prompt ───────────────────────────────────────────────────────────

def _confirm_mode() -> bool:
    """
    Ask the user whether to run in demo mode or full generation mode.
    Returns True if demo mode, False if full mode.
    Defaults to demo mode after 2 minutes.
    """
    import shutil
    has_claude = bool(shutil.which("claude"))

    if not has_claude:
        header(
            "LESSON CREATOR  —  Outwood Grange Academies Trust",
            "Demo mode  (claude CLI not found on PATH)",
        )
        notify("Running in demo mode — stub content will be used", "info")
        notify("Install Claude Code to generate real AI content", "info")
        return True  # no choice needed

    header(
        "LESSON CREATOR  —  Outwood Grange Academies Trust",
        "Claude Code detected — choose generation mode",
    )

    choice = prompt_choice(
        title="SELECT GENERATION MODE",
        options=[
            ("Demo mode",
             "Fast — uses realistic stub content, no AI calls, no cost"),
            ("Full generation (uses Claude Code)",
             "Slow — real AI-generated maths content via local claude CLI"),
            ("Cancel",
             "Exit without generating"),
        ],
        default=0,
        timeout=120,
    )

    if choice == 2:
        notify("Cancelled.", "warning")
        sys.exit(0)

    return choice == 0  # 0 = demo, 1 = full


# ── CLI entry point ──────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Generate an Outwood lesson slide deck from a Scheme of Work.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=(
            "Run with no arguments for the interactive module browser.\n"
            "Supply --summary and --methods to skip module selection."
        ),
    )
    parser.add_argument("--summary",
                        help="Path to 'Core Plus Summary of Intent' PDF")
    parser.add_argument("--methods",
                        help="Path to 'Common Methods / Consistent Methodology' PDF")
    parser.add_argument("--topic",    default="",
                        help="Topic name override (e.g. 'Algebra 1 Core Plus')")
    parser.add_argument("--year",     default="",
                        help="Year group (e.g. 'Year 9') — skips the interactive prompt")
    parser.add_argument("--pitch",    default="",
                        help="Teaching pitch (Nurture / Core / Core Plus / Extension)")
    parser.add_argument("--hours",    type=int, default=0,
                        help="Override recommended teaching hours")
    parser.add_argument("--output",   default="",
                        help="Output PPTX path (default: Output/<topic>.pptx)")
    args = parser.parse_args()

    # ── Mode selection ────────────────────────────────────────
    demo = _confirm_mode()
    content_gen.set_demo_mode(demo)

    mode_label = "DEMO" if demo else "FULL (API)"
    notify(f"Mode: {mode_label}", "info")

    # ── Module / PDF resolution ───────────────────────────────
    summary_path = args.summary
    methods_path = args.methods
    topic_override = args.topic
    module_tier = ""

    if not summary_path or not methods_path:
        result = module_selector.select_module()
        if result is None:
            # Module browser cancelled — ask for paths manually
            if not summary_path:
                summary_path = _ask_pdf_path("Summary of Intent")
            if not methods_path:
                methods_path = _ask_pdf_path("Common Methods")
            if not summary_path or not methods_path:
                notify("Cannot generate without both PDF files — exiting", "error")
                sys.exit(1)
        else:
            summary_path   = result["summary_pdf"]
            methods_path   = result["methods_pdf"]
            topic_override = topic_override or result["topic_name"]
            module_tier    = result["module"].get("tier_name", "")

    # ── Class profile ─────────────────────────────────────────
    profile = cp.select_profile(
        year=args.year,
        pitch=args.pitch,
        default_pitch=module_tier or "Core Plus",
    )
    content_gen.set_class_profile(profile)

    # ── Parse PDFs ────────────────────────────────────────────
    notify("Parsing Summary of Intent …", "progress")
    sow = sow_parser.parse_summary(summary_path)

    if topic_override:
        sow["topic"] = topic_override
    if args.hours:
        sow["time_hours"] = args.hours

    topic_name  = sow["topic"]
    output_path = args.output or f"Output/{topic_name.replace(' ', '_')}.pptx"

    notify("Parsing Common Methods …", "progress")
    methods = sow_parser.parse_methods(methods_path)

    notify(f"Topic:   {topic_name}", "info")
    notify(f"Class:   {profile['label']}", "info")
    notify(f"Output:  {output_path}", "info")

    build_deck(sow, methods, topic_name, output_path)


def _ask_pdf_path(label: str) -> str:
    """Fallback: ask the user to enter a PDF path at the terminal."""
    print(f"\n  Enter path to {label} PDF (or leave blank to cancel):")
    print("  > ", end="", flush=True)
    path = sys.stdin.readline().strip().strip('"').strip("'")
    if not path:
        return ""
    from pathlib import Path
    if not Path(path).exists():
        notify(f"File not found: {path}", "warning")
        return ""
    return path


if __name__ == "__main__":
    main()
