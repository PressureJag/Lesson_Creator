"""Build individual slides using python-pptx, applying the Outwood brand theme."""

import io
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from . import theme as T

# ── Layout geometry (16:9 slide: 13.33" × 7.5") ─────────────────────────────
_L   = Inches(0.25)    # left margin
_CT  = Inches(0.68)    # content top (below taller header)
_CW  = Inches(12.83)   # full content width
_CH  = Inches(6.57)    # usable content height
_COL = Inches(6.14)    # column width for 2-column layouts
_R2  = Inches(6.69)    # right-column x start  (L + COL + 0.30" gap)


# ── Slide / shape helpers ─────────────────────────────────────────────────────

def _new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = T.SLIDE_WIDTH
    prs.slide_height = T.SLIDE_HEIGHT
    return prs


def _blank_slide(prs: Presentation):
    """Blank slide with the Outwood cream-yellow background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = T.YELLOW_BG
    return slide


def _textbox(slide, left, top, width, height,
             text="", font_size=None, bold=False,
             colour=None, alignment=PP_ALIGN.LEFT,
             underline=False, italic=False, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = T.FONT_NAME
    run.font.size = font_size or T.BODY_SIZE
    run.font.bold = bold
    run.font.italic = italic
    run.font.underline = underline
    if colour:
        run.font.color.rgb = colour
    return txBox


def _box(slide, left, top, width, height, fill_colour,
         text="", text_colour=None, font_size=None, bold=False,
         alignment=PP_ALIGN.CENTER, border_colour=None, border_pt=1.5,
         word_wrap=True):
    """Filled rectangle with optional text and border."""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_colour
    if border_colour:
        shape.line.color.rgb = border_colour
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = word_wrap
        p = tf.paragraphs[0]
        p.alignment = alignment
        run = p.add_run()
        run.text = text
        run.font.name = T.FONT_NAME
        run.font.size = font_size or T.BODY_SIZE
        run.font.bold = bold
        run.font.color.rgb = text_colour or T.WHITE
    return shape


def _rounded_box(slide, left, top, width, height, fill_colour,
                 text="", text_colour=None, font_size=None, bold=False,
                 alignment=PP_ALIGN.CENTER, border_colour=None, border_pt=2.0,
                 word_wrap=True, rounding=0.06):
    """Rounded-rectangle card (shape type 5) with optional text."""
    shape = slide.shapes.add_shape(5, left, top, width, height)  # 5 = ROUNDED_RECTANGLE
    shape.adjustments[0] = rounding
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_colour
    if border_colour:
        shape.line.color.rgb = border_colour
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = word_wrap
        p = tf.paragraphs[0]
        p.alignment = alignment
        run = p.add_run()
        run.text = text
        run.font.name = T.FONT_NAME
        run.font.size = font_size or T.BODY_SIZE
        run.font.bold = bold
        run.font.color.rgb = text_colour or T.WHITE
    return shape


def _circle(slide, cx, cy, diameter, fill_colour, text="",
            text_colour=None, font_size=None, bold=True):
    """Filled circle with centred text (used for step-number badges)."""
    r = diameter / 2
    shape = slide.shapes.add_shape(9, cx - r, cy - r, diameter, diameter)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_colour
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.name = T.FONT_NAME
        run.font.size = font_size or Pt(18)
        run.font.bold = bold
        run.font.color.rgb = text_colour or T.WHITE
    return shape


def _embed_image(slide, img_bytes: io.BytesIO, left, top, width, height):
    slide.shapes.add_picture(img_bytes, left, top, width, height)


# ── Header system ─────────────────────────────────────────────────────────────

def _add_header(slide, slide_type: str, subtitle: str = "",
                badge_text: str = "", badge_colour=None) -> None:
    """
    Three-element header used on every slide:
      1. OGAT badge     — small cream square, navy border, top-left
      2. Header bar     — rounded rect, navy border, slide_type + subtitle text
      3. Phase badge    — coloured circle, top-right  (omitted when badge_text is empty)
    """
    H_TOP   = Inches(0.06)
    H_H     = Inches(0.54)
    BADGE_W = Inches(0.90)
    BADGE_H = Inches(0.58)

    # 1 ── OGAT badge
    ogat = slide.shapes.add_shape(1,
        Inches(0.07), H_TOP, Inches(0.54), H_H)
    ogat.fill.solid()
    ogat.fill.fore_color.rgb = T.YELLOW_BG
    ogat.line.color.rgb = T.NAVY
    ogat.line.width = Pt(2.0)
    tf = ogat.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "OGAT"
    run.font.name = T.FONT_NAME
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = T.NAVY

    # 2 ── Header bar
    badge_space = (BADGE_W + Inches(0.10)) if badge_text else Inches(0.05)
    hbar_left  = Inches(0.66)
    hbar_width = T.SLIDE_WIDTH - hbar_left - badge_space

    hbar = slide.shapes.add_shape(1,
        hbar_left, H_TOP, hbar_width, H_H)
    hbar.fill.solid()
    hbar.fill.fore_color.rgb = T.WHITE
    hbar.line.color.rgb = T.NAVY
    hbar.line.width = Pt(2.0)

    tf = hbar.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    run1 = p.add_run()
    run1.text = slide_type
    run1.font.name = T.FONT_NAME
    run1.font.size = Pt(26)
    run1.font.bold = True
    run1.font.color.rgb = T.NAVY

    if subtitle:
        run2 = p.add_run()
        run2.text = f"   {subtitle}"
        run2.font.name = T.FONT_NAME
        run2.font.size = Pt(14)
        run2.font.bold = False
        run2.font.color.rgb = T.NAVY

    # 3 ── Phase badge (circle)
    if badge_text and badge_colour:
        badge_left = T.SLIDE_WIDTH - BADGE_W - Inches(0.05)
        badge = slide.shapes.add_shape(9,   # oval
            badge_left, Inches(0.02), BADGE_W, BADGE_H)
        badge.fill.solid()
        badge.fill.fore_color.rgb = badge_colour
        badge.line.fill.background()
        tf = badge.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = badge_text
        run.font.name = T.FONT_NAME
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = T.WHITE


# ── Slide constructors ────────────────────────────────────────────────────────

def make_title_slide(prs: Presentation, topic_name: str,
                     big_question: str, time_hours: int) -> None:
    """Opening title card — navy strips top & bottom, centred content card."""
    slide = _blank_slide(prs)

    # Navy strips
    _box(slide, Inches(0), Inches(0),
         T.SLIDE_WIDTH, Inches(0.55), T.NAVY)
    _box(slide, Inches(0), T.SLIDE_HEIGHT - Inches(0.55),
         T.SLIDE_WIDTH, Inches(0.55), T.NAVY)

    # Central white card
    _box(slide, Inches(0.5), Inches(0.65),
         Inches(12.33), Inches(6.25),
         T.WHITE, border_colour=T.NAVY, border_pt=2.0)

    # Purple pill banner inside card
    _box(slide, Inches(2.5), Inches(0.90),
         Inches(8.33), Inches(0.60),
         T.NAVY,
         text=topic_name.upper(),
         font_size=Pt(16), bold=True, alignment=PP_ALIGN.CENTER)

    # Deck title
    _textbox(slide, Inches(1.0), Inches(1.65),
             Inches(11.33), Inches(1.20),
             text="Full sequence",
             font_size=Pt(40), bold=True, colour=T.NAVY,
             alignment=PP_ALIGN.CENTER)

    # Big question / subtitle
    _textbox(slide, Inches(1.0), Inches(2.95),
             Inches(11.33), Inches(0.70),
             text=big_question,
             font_size=Pt(18), bold=True, colour=T.PINK,
             alignment=PP_ALIGN.CENTER)

    # Recommended time
    _textbox(slide, Inches(1.0), Inches(3.75),
             Inches(11.33), Inches(0.50),
             text=f"Recommended {time_hours} hours",
             font_size=Pt(14), colour=T.DARK_GREY,
             alignment=PP_ALIGN.CENTER)

    # Footer
    _textbox(slide, Inches(2.0), Inches(6.90),
             Inches(9.33), Inches(0.40),
             text="OGAT Maths  ·  Core Plus Tier",
             font_size=Pt(12), colour=T.DARK_GREY,
             alignment=PP_ALIGN.CENTER)


def make_overview_slide(prs: Presentation, topic_name: str,
                        prior_points: list[str], objectives: list[str],
                        future_points: list[str]) -> None:
    """Big Picture — three-column map showing where this module sits in the learning journey."""
    slide = _blank_slide(prs)
    _add_header(slide, "Big Picture",
                f"Where {topic_name} fits in your learning journey")

    col_top = _CT + Inches(0.15)
    col_h   = _CH - Inches(0.20)
    side_w  = Inches(2.60)
    mid_w   = Inches(6.40)
    arrow_w = Inches(0.40)

    left_x  = _L
    arr1_x  = left_x + side_w
    mid_x   = arr1_x + arrow_w
    arr2_x  = mid_x + mid_w
    right_x = arr2_x + arrow_w
    arrow_y = col_top + Inches(1.5)

    # Left — Prior Knowledge
    _box(slide, left_x, col_top, side_w, col_h,
         T.CARD_PEACH, border_colour=T.NAVY, border_pt=1.5)
    _textbox(slide, left_x + Inches(0.12), col_top + Inches(0.10),
             side_w - Inches(0.24), Inches(0.40),
             text="Prior Knowledge", font_size=Pt(13), bold=True, colour=T.NAVY)
    prior_text = "\n".join(f"• {p}" for p in prior_points[:5])
    _textbox(slide, left_x + Inches(0.12), col_top + Inches(0.55),
             side_w - Inches(0.24), col_h - Inches(0.65),
             text=prior_text, font_size=Pt(10), colour=T.NAVY, word_wrap=True)

    # Arrows
    _textbox(slide, arr1_x, arrow_y, arrow_w, Inches(0.35),
             text="→", font_size=Pt(22), bold=True, colour=T.NAVY,
             alignment=PP_ALIGN.CENTER)
    _textbox(slide, arr2_x, arrow_y, arrow_w, Inches(0.35),
             text="→", font_size=Pt(22), bold=True, colour=T.NAVY,
             alignment=PP_ALIGN.CENTER)

    # Middle — Current Module objectives
    _box(slide, mid_x, col_top, mid_w, col_h,
         T.YELLOW_BG, border_colour=T.NAVY, border_pt=2.0)
    _textbox(slide, mid_x + Inches(0.15), col_top + Inches(0.10),
             mid_w - Inches(0.30), Inches(0.50),
             text=topic_name, font_size=Pt(15), bold=True,
             colour=T.NAVY, alignment=PP_ALIGN.CENTER)
    obj_text = "\n".join(f"{i + 1}.  {o}" for i, o in enumerate(objectives))
    _textbox(slide, mid_x + Inches(0.15), col_top + Inches(0.72),
             mid_w - Inches(0.30), col_h - Inches(0.82),
             text=obj_text, font_size=Pt(11), colour=T.NAVY, word_wrap=True)

    # Right — Future Learning
    _box(slide, right_x, col_top, side_w, col_h,
         T.CARD_GREEN, border_colour=T.NAVY, border_pt=1.5)
    _textbox(slide, right_x + Inches(0.12), col_top + Inches(0.10),
             side_w - Inches(0.24), Inches(0.40),
             text="Future Learning", font_size=Pt(13), bold=True, colour=T.NAVY)
    future_text = "\n".join(f"• {f}" for f in future_points[:5])
    _textbox(slide, right_x + Inches(0.12), col_top + Inches(0.55),
             side_w - Inches(0.24), col_h - Inches(0.65),
             text=future_text, font_size=Pt(10), colour=T.NAVY, word_wrap=True)


def make_prior_knowledge(prs: Presentation, topic_name: str,
                         prior: list[str]) -> None:
    """Prior knowledge recap slide."""
    slide = _blank_slide(prs)
    _add_header(slide, "Prior Knowledge — What do you already know?")

    text = "\n".join(f"•  {p}" for p in prior)
    _textbox(slide, _L, _CT + Inches(0.15),
             _CW, _CH - Inches(0.15),
             text=text, font_size=Pt(17), colour=T.BLACK, word_wrap=True)


def make_vocabulary(prs: Presentation, topic_name: str,
                    vocab_list: list[str]) -> None:
    """Key vocabulary slide."""
    slide = _blank_slide(prs)
    _add_header(slide, "Key Vocabulary")

    cols = 3
    per_col = (len(vocab_list) + cols - 1) // cols
    col_w = Inches(4.1)
    for ci in range(cols):
        chunk = vocab_list[ci * per_col: (ci + 1) * per_col]
        text = "\n".join(f"•  {v}" for v in chunk)
        _textbox(slide,
                 _L + ci * (col_w + Inches(0.15)),
                 _CT + Inches(0.15),
                 col_w, _CH,
                 text=text, font_size=Pt(16),
                 colour=T.BLACK, word_wrap=True)


def make_starter_plus(prs: Presentation, topic_name: str,
                      questions: list[str], answers: list[str] = None) -> None:
    """4-question starter in a 2×2 rounded-card grid."""
    is_answers = bool(answers)
    slide = _blank_slide(prs)

    if is_answers:
        _add_header(slide, "STARTER", "let's check our answers",
                    "Feed-back", T.TEAL)
    else:
        _add_header(slide, "STARTER", "Can you remember what we've learned?",
                    "Recap & Recall", T.ORANGE)

    card_colours = [T.CARD_YELLOW, T.CARD_PEACH, T.CARD_GREEN, T.CARD_BLUE]
    # Slightly darker border per card (gives a subtle edge)
    border_colours = [
        RGBColor(0xD4, 0xC0, 0x50),
        RGBColor(0xD4, 0xA0, 0x60),
        RGBColor(0x70, 0xB8, 0x80),
        RGBColor(0x70, 0xAA, 0xD0),
    ]
    card_icons = ["📝", "💡", "🎯", "⭐"]

    CARD_W = _COL
    CARD_H = Inches(3.10)
    GAP    = Inches(0.15)

    positions = [
        (_L,  _CT + Inches(0.08)),
        (_R2, _CT + Inches(0.08)),
        (_L,  _CT + Inches(0.08) + CARD_H + GAP),
        (_R2, _CT + Inches(0.08) + CARD_H + GAP),
    ]

    q_list = (list(questions) + [""] * 4)[:4]
    a_list = (list(answers) if answers else []) + [""] * 4
    answer_colours = [
        RGBColor(0xC0, 0x39, 0x2B),  # deep red
        RGBColor(0x1E, 0x5B, 0xB8),  # deep blue
        RGBColor(0x1E, 0x7A, 0x3C),  # deep green
        RGBColor(0x6A, 0x1B, 0x9A),  # deep purple
    ]

    for i, ((lft, top), colour, border) in enumerate(
            zip(positions, card_colours, border_colours)):

        _rounded_box(slide, lft, top, CARD_W, CARD_H,
                     colour, border_colour=border, border_pt=2.0, rounding=0.06)

        # Question number (bold, large)
        _textbox(slide, lft + Inches(0.14), top + Inches(0.12),
                 Inches(0.45), Inches(0.42),
                 text=f"{i + 1}.", font_size=Pt(20), bold=True, colour=T.NAVY)

        # Emoji icon — top-right of card
        _textbox(slide, lft + CARD_W - Inches(0.55), top + Inches(0.08),
                 Inches(0.48), Inches(0.40),
                 text=card_icons[i], font_size=Pt(20),
                 alignment=PP_ALIGN.CENTER)

        body = q_list[i]
        if is_answers and i < len(answers) and a_list[i]:
            # Question text (compact, top of card)
            _textbox(slide, lft + Inches(0.55), top + Inches(0.10),
                     CARD_W - Inches(1.10), Inches(0.80),
                     text=body, font_size=Pt(16), colour=T.NAVY, word_wrap=True)
            # Answer in coloured bold text
            _textbox(slide, lft + Inches(0.16), top + Inches(1.00),
                     CARD_W - Inches(0.32), CARD_H - Inches(1.20),
                     text=a_list[i], font_size=Pt(16), bold=True,
                     colour=answer_colours[i], word_wrap=True)
            # Green checkmark bottom-right
            _textbox(slide, lft + CARD_W - Inches(0.48), top + CARD_H - Inches(0.42),
                     Inches(0.40), Inches(0.38),
                     text="✓", font_size=Pt(20), bold=True,
                     colour=RGBColor(0x1E, 0x8A, 0x1E),
                     alignment=PP_ALIGN.CENTER)
        else:
            # Question text
            _textbox(slide, lft + Inches(0.55), top + Inches(0.10),
                     CARD_W - Inches(1.10), Inches(0.95),
                     text=body, font_size=Pt(17), colour=T.NAVY, word_wrap=True)
            # "Working ..." prompt
            _textbox(slide, lft + Inches(0.16), top + Inches(1.12),
                     CARD_W - Inches(0.32), Inches(0.28),
                     text="Working ...", font_size=Pt(12), italic=True,
                     colour=RGBColor(0xC8, 0x7A, 0x00))
            # Working lines (subtle dashes)
            for line_n in range(3):
                line_y = top + Inches(1.52) + line_n * Inches(0.50)
                _box(slide, lft + Inches(0.16), line_y,
                     CARD_W - Inches(0.32), Inches(0.025),
                     RGBColor(0xBB, 0xBB, 0xBB))


def make_learning_intro(prs: Presentation, topic_name: str,
                        objective: str, objective_num: int) -> None:
    """Clarity of Learning — objective intro slide."""
    slide = _blank_slide(prs)
    _add_header(slide,
                f"Clarity of Learning — Objective {objective_num}",
                "", "Clarity", T.PINK)

    # Topic Question box (pink border)
    tq_h = Inches(0.90)
    _box(slide, _L, _CT + Inches(0.10),
         _CW, tq_h,
         T.WHITE, border_colour=T.PINK, border_pt=2.0)
    _textbox(slide, _L + Inches(0.15), _CT + Inches(0.10),
             _CW - Inches(0.30), tq_h,
             text="Topic Question", font_size=Pt(13), bold=True,
             colour=T.PINK)
    _textbox(slide, _L + Inches(0.15), _CT + Inches(0.38),
             _CW - Inches(0.30), Inches(0.55),
             text=f"How do I {objective.lower().rstrip('.')}?",
             font_size=Pt(14), colour=T.NAVY)

    # Lesson Question box (teal border)
    lq_top = _CT + Inches(1.10)
    lq_h   = Inches(0.90)
    _box(slide, _L, lq_top,
         _CW, lq_h,
         T.WHITE, border_colour=T.TEAL, border_pt=2.0)
    _textbox(slide, _L + Inches(0.15), lq_top,
             _CW - Inches(0.30), lq_h,
             text="Lesson Question", font_size=Pt(13), bold=True,
             colour=T.TEAL)
    _textbox(slide, _L + Inches(0.15), lq_top + Inches(0.28),
             _CW - Inches(0.30), Inches(0.55),
             text=objective,
             font_size=Pt(14), colour=T.NAVY)

    # "What it'll look like" area (dark border)
    wil_top = _CT + Inches(2.10)
    wil_h   = Inches(2.70)
    _box(slide, _L, wil_top,
         _CW, wil_h,
         T.WHITE, border_colour=T.NAVY, border_pt=1.5)
    _textbox(slide, _L + Inches(0.15), wil_top + Inches(0.10),
             _CW - Inches(0.30), Inches(0.35),
             text="What it'll look like ...",
             font_size=Pt(13), bold=True, colour=T.NAVY)

    # Success criteria bar
    sc_top = _CT + Inches(4.90)
    sc_h   = Inches(1.30)
    _box(slide, _L, sc_top, _CW, sc_h,
         T.YELLOW_BG, border_colour=T.NAVY, border_pt=1.0)
    _textbox(slide, _L + Inches(0.15), sc_top + Inches(0.08),
             _CW - Inches(0.30), Inches(0.30),
             text="Success criteria  —  I can ...",
             font_size=Pt(12), bold=True, colour=T.NAVY)

    # 4 success-criteria chips
    chip_colours = [T.CARD_GREEN, T.CARD_PEACH, T.CARD_BLUE, T.PURPLE_LIGHT]
    chip_w = Inches(2.90)
    chip_h = Inches(0.55)
    chip_y = sc_top + Inches(0.55)
    for ci, chip_col in enumerate(chip_colours):
        chip_x = _L + Inches(0.15) + ci * (chip_w + Inches(0.20))
        _box(slide, chip_x, chip_y, chip_w, chip_h,
             chip_col, border_colour=chip_col, border_pt=1.0)


def make_ido_slide(prs: Presentation, topic_name: str,
                   heading: str, worked_example: str,
                   method_text: str = "", notes: str = "",
                   diagram: "io.BytesIO | None" = None) -> None:
    """I Do — teacher-led worked example with numbered step circles.

    diagram: optional BytesIO PNG from diagram_gen — shown in the left panel
             instead of (or above) the method text.
    """
    slide = _blank_slide(prs)
    _add_header(slide, "Worked Example", heading, "I do", T.PURPLE)

    # Question banner (yellow, rounded)
    _rounded_box(slide, _L, _CT + Inches(0.08),
                 _CW, Inches(0.62),
                 T.CARD_YELLOW, text=heading,
                 text_colour=T.NAVY, font_size=Pt(20), bold=True,
                 alignment=PP_ALIGN.CENTER, rounding=0.08)

    content_top = _CT + Inches(0.82)
    content_h   = Inches(4.35)
    method_w    = Inches(4.10)
    gap         = Inches(0.15)
    right_x     = _L + method_w + gap
    right_w     = _CW - method_w - gap

    if method_text or diagram:
        # Left — Common Method panel (rounded, peach/orange)
        _rounded_box(slide, _L, content_top, method_w, content_h,
                     T.CARD_PEACH, border_colour=T.ORANGE, border_pt=2.5,
                     rounding=0.06)
        _textbox(slide, _L + Inches(0.14), content_top + Inches(0.10),
                 method_w - Inches(0.28), Inches(0.32),
                 text="COMMON METHOD", font_size=Pt(12), bold=True,
                 colour=T.ORANGE)

        label_bottom = content_top + Inches(0.46)

        if diagram:
            # Embed diagram image, preserving natural aspect ratio within the panel
            diag_w = method_w - Inches(0.24)
            diag_h = Inches(2.20)
            _embed_image(slide, diagram,
                         _L + Inches(0.12), label_bottom,
                         diag_w, diag_h)
            text_top = label_bottom + diag_h + Inches(0.10)
            text_h   = content_top + content_h - text_top - Inches(0.10)
            if method_text and text_h > Inches(0.30):
                _textbox(slide, _L + Inches(0.14), text_top,
                         method_w - Inches(0.28), text_h,
                         text=method_text, font_size=Pt(11),
                         colour=T.NAVY, word_wrap=True)
        else:
            _textbox(slide, _L + Inches(0.14), label_bottom,
                     method_w - Inches(0.28), content_h - Inches(0.56),
                     text=method_text, font_size=Pt(12),
                     colour=T.NAVY, word_wrap=True)

        # Right — Worked example panel (rounded, white)
        _rounded_box(slide, right_x, content_top,
                     right_w, content_h,
                     T.WHITE, border_colour=T.NAVY, border_pt=1.5,
                     rounding=0.06)

        # Parse the worked example into numbered steps using circle badges.
        # Spacing is calibrated so the full sequence fits in 4.19" with Comic Sans:
        #   body 0.30" · blank 0.08" · step 0.48" · in-general 0.52" → total ≈4.16"
        lines = worked_example.strip().split("\n")
        step_colours = [
            RGBColor(0xF5, 0x9A, 0x1C),  # orange  — Step 1
            RGBColor(0x17, 0xA5, 0x89),  # teal    — Step 2
            RGBColor(0x5B, 0x2C, 0x8D),  # purple  — Step 3
            RGBColor(0xE5, 0x00, 0x7D),  # pink    — Step 4
        ]
        step_num = 0
        y_cursor = content_top + Inches(0.16)
        circ_d   = Inches(0.36)
        line_h   = Inches(0.26)   # reduced from 0.30 to keep full sequence in panel

        for line in lines:
            stripped = line.strip()
            if not stripped:
                y_cursor += Inches(0.08)   # tighter blank-line gap
                continue
            is_step = stripped.lower().startswith("step ")
            if is_step and step_num < len(step_colours):
                col = step_colours[step_num]
                cx = right_x + Inches(0.30)
                cy = y_cursor + circ_d / 2
                _circle(slide, cx, cy, circ_d, col,
                        text=str(step_num + 1), font_size=Pt(14))
                _textbox(slide, right_x + Inches(0.72), y_cursor,
                         right_w - Inches(0.84), line_h + Inches(0.10),
                         text=stripped, font_size=Pt(14), bold=True,
                         colour=col, word_wrap=True)
                step_num += 1
                y_cursor += line_h + Inches(0.12)
            elif stripped.startswith("[") or stripped.startswith("("):
                _textbox(slide, right_x + Inches(0.72), y_cursor,
                         right_w - Inches(0.84), line_h,
                         text=stripped, font_size=Pt(12), italic=True,
                         colour=T.DARK_GREY, word_wrap=True)
                y_cursor += Inches(0.22)
            elif stripped.lower().startswith("in general"):
                _rounded_box(slide, right_x + Inches(0.14),
                             y_cursor + Inches(0.04),
                             right_w - Inches(0.28), line_h + Inches(0.10),
                             RGBColor(0xE8, 0xF8, 0xE8),
                             text=stripped, text_colour=RGBColor(0x1E, 0x7A, 0x3C),
                             font_size=Pt(13), bold=True,
                             alignment=PP_ALIGN.LEFT, rounding=0.08)
                y_cursor += line_h + Inches(0.18)
            else:
                _textbox(slide, right_x + Inches(0.18), y_cursor,
                         right_w - Inches(0.30), line_h,
                         text=stripped, font_size=Pt(13),
                         colour=T.NAVY, word_wrap=True)
                y_cursor += line_h + Inches(0.04)
    else:
        _rounded_box(slide, _L, content_top, _CW, content_h,
                     T.WHITE, border_colour=T.NAVY, border_pt=1.5,
                     rounding=0.06)
        _textbox(slide, _L + Inches(0.18), content_top + Inches(0.15),
                 _CW - Inches(0.36), content_h - Inches(0.30),
                 text=worked_example, font_size=Pt(16),
                 colour=T.NAVY, word_wrap=True)

    # Green answer bar (rounded)
    ans_top = _CT + Inches(5.25)
    _rounded_box(slide, _L, ans_top, _CW, Inches(0.52),
                 RGBColor(0x1E, 0x7A, 0x4A),
                 text="Answer shown above — check your working",
                 text_colour=T.WHITE, font_size=Pt(15), bold=True,
                 rounding=0.08)

    # Navy key-idea bar
    key_idea = notes if notes else "Key idea:  show every step — method first, then substitute"
    _rounded_box(slide, _L, ans_top + Inches(0.57), _CW, Inches(0.44),
                 T.NAVY, text=key_idea,
                 text_colour=T.CARD_YELLOW, font_size=Pt(13), bold=True,
                 rounding=0.08)


def make_wedo_slide(prs: Presentation, topic_name: str,
                    heading: str, question: str,
                    scaffold_steps: list[str] = None) -> None:
    """We Do — 4-card guided practice with teacher note bar."""
    slide = _blank_slide(prs)
    _add_header(slide, "We do", heading, "We do", T.TEAL)

    # Teacher note bar (lavender strip below header)
    note_h = Inches(0.30)
    _rounded_box(slide, _L, _CT + Inches(0.04),
                 _CW, note_h,
                 T.PURPLE_LIGHT,
                 text=question,
                 text_colour=T.NAVY, font_size=Pt(12),
                 alignment=PP_ALIGN.LEFT, rounding=0.05)

    card_colours = [T.CARD_YELLOW, T.CARD_PEACH, T.CARD_GREEN, T.CARD_BLUE]
    border_colours = [
        RGBColor(0xD4, 0xC0, 0x50),
        RGBColor(0xD4, 0xA0, 0x60),
        RGBColor(0x70, 0xB8, 0x80),
        RGBColor(0x70, 0xAA, 0xD0),
    ]
    card_icons = ["📝", "💡", "🎯", "⭐"]

    CARD_W   = _COL
    CARD_H   = Inches(2.95)
    GAP      = Inches(0.15)
    cards_top = _CT + Inches(0.42)

    positions = [
        (_L,  cards_top),
        (_R2, cards_top),
        (_L,  cards_top + CARD_H + GAP),
        (_R2, cards_top + CARD_H + GAP),
    ]

    steps = (scaffold_steps or [
        "Try it — show your working step by step.",
        "Check — does your answer make sense?",
        "Write the answer clearly.",
        "Can you spot a pattern?",
    ])[:4]

    for i, ((lft, top), colour, border) in enumerate(
            zip(positions, card_colours, border_colours)):

        _rounded_box(slide, lft, top, CARD_W, CARD_H,
                     colour, border_colour=border, border_pt=2.0, rounding=0.06)

        _textbox(slide, lft + Inches(0.14), top + Inches(0.10),
                 Inches(0.45), Inches(0.42),
                 text=f"{i + 1}.", font_size=Pt(20), bold=True, colour=T.NAVY)

        _textbox(slide, lft + CARD_W - Inches(0.55), top + Inches(0.08),
                 Inches(0.48), Inches(0.40),
                 text=card_icons[i], font_size=Pt(20),
                 alignment=PP_ALIGN.CENTER)

        step = steps[i] if i < len(steps) else ""
        _textbox(slide, lft + Inches(0.55), top + Inches(0.10),
                 CARD_W - Inches(1.10), Inches(0.90),
                 text=step, font_size=Pt(15), colour=T.NAVY, word_wrap=True)

        _textbox(slide, lft + Inches(0.16), top + Inches(1.08),
                 CARD_W - Inches(0.32), Inches(0.28),
                 text="Working ...", font_size=Pt(12), italic=True,
                 colour=RGBColor(0xC8, 0x7A, 0x00))

        for line_n in range(3):
            line_y = top + Inches(1.46) + line_n * Inches(0.48)
            _box(slide, lft + Inches(0.16), line_y,
                 CARD_W - Inches(0.32), Inches(0.025),
                 RGBColor(0xBB, 0xBB, 0xBB))


def make_youdo_slide(prs: Presentation, topic_name: str,
                     heading: str, question: str,
                     answer: str = "") -> None:
    """You Do — 4-card independent practice."""
    slide = _blank_slide(prs)
    _add_header(slide, "You do", heading, "You do", T.PURPLE)

    note_h = Inches(0.30)
    _rounded_box(slide, _L, _CT + Inches(0.04),
                 _CW, note_h,
                 T.PURPLE_LIGHT,
                 text=question,
                 text_colour=T.NAVY, font_size=Pt(12),
                 alignment=PP_ALIGN.LEFT, rounding=0.05)

    card_colours = [T.CARD_YELLOW, T.CARD_PEACH, T.CARD_GREEN, T.CARD_BLUE]
    border_colours = [
        RGBColor(0xD4, 0xC0, 0x50),
        RGBColor(0xD4, 0xA0, 0x60),
        RGBColor(0x70, 0xB8, 0x80),
        RGBColor(0x70, 0xAA, 0xD0),
    ]
    card_icons = ["📝", "💡", "🎯", "⭐"]

    CARD_W    = _COL
    CARD_H    = Inches(2.95)
    GAP       = Inches(0.15)
    cards_top = _CT + Inches(0.42)

    positions = [
        (_L,  cards_top),
        (_R2, cards_top),
        (_L,  cards_top + CARD_H + GAP),
        (_R2, cards_top + CARD_H + GAP),
    ]
    labels = [
        "Try this question on your own.",
        "Show all working clearly.",
        "Write your final answer.",
        "Check — does it make sense?",
    ]

    for i, ((lft, top), colour, border) in enumerate(
            zip(positions, card_colours, border_colours)):

        _rounded_box(slide, lft, top, CARD_W, CARD_H,
                     colour, border_colour=border, border_pt=2.0, rounding=0.06)

        _textbox(slide, lft + Inches(0.14), top + Inches(0.10),
                 Inches(0.45), Inches(0.42),
                 text=f"{i + 1}.", font_size=Pt(20), bold=True, colour=T.NAVY)

        _textbox(slide, lft + CARD_W - Inches(0.55), top + Inches(0.08),
                 Inches(0.48), Inches(0.40),
                 text=card_icons[i], font_size=Pt(20),
                 alignment=PP_ALIGN.CENTER)

        _textbox(slide, lft + Inches(0.55), top + Inches(0.10),
                 CARD_W - Inches(1.10), Inches(0.90),
                 text=labels[i], font_size=Pt(15), colour=T.NAVY, word_wrap=True)

        _textbox(slide, lft + Inches(0.16), top + Inches(1.08),
                 CARD_W - Inches(0.32), Inches(0.28),
                 text="Working ...", font_size=Pt(12), italic=True,
                 colour=RGBColor(0xC8, 0x7A, 0x00))

        for line_n in range(3):
            line_y = top + Inches(1.46) + line_n * Inches(0.48)
            _box(slide, lft + Inches(0.16), line_y,
                 CARD_W - Inches(0.32), Inches(0.025),
                 RGBColor(0xBB, 0xBB, 0xBB))

    if answer:
        _textbox(slide, _L, _CT + _CH - Inches(0.35),
                 _CW, Inches(0.30),
                 text=f"Answer (teacher reference): {answer}",
                 font_size=Pt(11), colour=T.DARK_GREY,
                 italic=True, alignment=PP_ALIGN.RIGHT)


def make_mini_whiteboard(prs: Presentation, topic_name: str,
                         question: str, question_num: int = None,
                         total: int = 10) -> None:
    """Mini whiteboard — one large question with display area."""
    slide = _blank_slide(prs)
    counter = f"  {question_num}/{total}" if question_num else ""
    _add_header(slide, f"Mini Whiteboard{counter}")

    _rounded_box(slide, _L, _CT + Inches(0.12),
                 _CW, Inches(5.85),
                 T.WHITE, text=question,
                 text_colour=T.NAVY, font_size=Pt(30),
                 alignment=PP_ALIGN.CENTER,
                 border_colour=T.TEAL, border_pt=3.0, rounding=0.05)

    _textbox(slide, _L, _CT + Inches(6.10), _CW, Inches(0.35),
             text="Show me your answer on your mini whiteboard.",
             font_size=Pt(13), colour=T.DARK_GREY, italic=True,
             alignment=PP_ALIGN.CENTER)


def make_independent_practice(prs: Presentation, topic_name: str,
                               questions: list[str]) -> None:
    """Independent practice — up to 10 numbered questions in two columns."""
    slide = _blank_slide(prs)
    _add_header(slide, "Independent Practice", "", "Indep.", T.ORANGE)

    qs  = list(questions)[:10]
    mid = (len(qs) + 1) // 2
    left_qs  = qs[:mid]
    right_qs = qs[mid:]

    left_text = "\n\n".join(f"{i + 1}.  {q}" for i, q in enumerate(left_qs))
    _textbox(slide, _L, _CT + Inches(0.10), _COL, _CH - Inches(0.10),
             text=left_text, font_size=Pt(14), colour=T.BLACK, word_wrap=True)

    if right_qs:
        right_text = "\n\n".join(
            f"{i + mid + 1}.  {q}" for i, q in enumerate(right_qs)
        )
        _textbox(slide, _R2, _CT + Inches(0.10), _COL, _CH - Inches(0.10),
                 text=right_text, font_size=Pt(14), colour=T.BLACK, word_wrap=True)
        _box(slide, _R2 - Inches(0.20), _CT, Inches(0.04), _CH,
             T.PURPLE_LIGHT)


def make_independent_answers(prs: Presentation, topic_name: str,
                              questions: list[str],
                              answers: list[str]) -> None:
    """Answers slide for independent practice."""
    slide = _blank_slide(prs)
    _add_header(slide, "Independent Practice", "let's check our answers",
                "Feed-back", T.TEAL)

    qs  = list(questions)[:10]
    ans = (list(answers) + ["?"] * 10)[:len(qs)]
    mid = (len(qs) + 1) // 2

    left_text = "\n\n".join(f"{i + 1}.  {ans[i]}" for i in range(min(mid, len(ans))))
    _textbox(slide, _L, _CT + Inches(0.10), _COL, _CH - Inches(0.10),
             text=left_text, font_size=Pt(14), colour=T.BLACK, word_wrap=True)

    if mid < len(qs):
        right_text = "\n\n".join(
            f"{i + mid + 1}.  {ans[i + mid]}" for i in range(len(qs) - mid)
        )
        _textbox(slide, _R2, _CT + Inches(0.10), _COL, _CH - Inches(0.10),
                 text=right_text, font_size=Pt(14), colour=T.BLACK, word_wrap=True)
        _box(slide, _R2 - Inches(0.20), _CT, Inches(0.04), _CH,
             T.PURPLE_LIGHT)


def make_plenary(prs: Presentation, topic_name: str,
                 objective: str, summary: str,
                 question: str, answer: str = "") -> None:
    """Plenary — lesson summary and final whiteboard check."""
    slide = _blank_slide(prs)
    _add_header(slide, "Plenary", "", "Plenary", T.PURPLE)

    _rounded_box(slide, _L, _CT + Inches(0.10),
                 _CW, Inches(1.65),
                 T.PURPLE_LIGHT, text=summary,
                 text_colour=T.NAVY, font_size=Pt(15),
                 alignment=PP_ALIGN.LEFT, rounding=0.06)

    _textbox(slide, _L, _CT + Inches(1.90),
             _CW, Inches(0.40),
             text="Now show me on your whiteboard:",
             font_size=Pt(16), bold=True, colour=T.NAVY)

    _rounded_box(slide, _L, _CT + Inches(2.40),
                 _CW, Inches(3.45),
                 T.WHITE, text=question,
                 text_colour=T.NAVY, font_size=Pt(26),
                 alignment=PP_ALIGN.CENTER,
                 border_colour=T.PURPLE, border_pt=2.5, rounding=0.06)

    if answer:
        _textbox(slide, _L, _CT + Inches(5.95),
                 _CW, Inches(0.45),
                 text=f"Answer: {answer}",
                 font_size=Pt(12), colour=T.DARK_GREY, italic=True)


def make_extension(prs: Presentation, topic_name: str,
                   heading: str, tasks: list[str]) -> None:
    """Extension task slide."""
    slide = _blank_slide(prs)
    _add_header(slide, "Extension", heading)

    text = "\n\n".join(f"•  {t}" for t in tasks)
    _box(slide, _L, _CT + Inches(0.15),
         _CW, _CH - Inches(0.15),
         RGBColor(0xFF, 0xF8, 0xE7),
         text=text, text_colour=T.BLACK,
         font_size=Pt(16), alignment=PP_ALIGN.LEFT)
